"""
Universal document extraction layer — Document Intelligence System.

Supported formats:
  .pdf          → PyMuPDF → pdfplumber → OCR (3-tier cascade, see pdf_utils.py)
  .docx         → python-docx  (paragraphs + tables)
  .pptx         → python-pptx  (one page per slide)
  .xlsx / .xls  → openpyxl + pandas (one page per sheet, sheets as tables)
  .csv          → stdlib csv   (content + table view)
  .txt / .md    → plain read   (50-line pages)
  .html / .htm  → BeautifulSoup (strip markup, treat as text)
  .jpg/.jpeg/.png/.bmp/.tiff → pytesseract OCR directly on image

All extractors return a DocumentContent (alias for PDFContent) so the rest
of the pipeline (chunker, vectorstore, document_service) is unchanged.
"""

import csv
import hashlib
import logging
import re
from pathlib import Path

# Re-export core types so callers only import from doc_utils
from app.utils.pdf_utils import (
    PDFContent as DocumentContent,  # noqa: F401 - public alias
    PageContent,
    TableData,
    PDFValidationError as DocumentExtractionError,
    extract_pdf,
    clean_text,
)
from app.config.settings import get_settings

settings = get_settings()
logger = logging.getLogger(__name__)

# ---------------------------------------------------------------------------
# Supported extensions
# ---------------------------------------------------------------------------

SUPPORTED_EXTENSIONS: set[str] = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx", ".xls",
    ".csv",
    ".txt", ".md",
    ".html", ".htm",
    ".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif",
}

# Human-readable label shown in the UI / error messages
FORMAT_LABELS: dict[str, str] = {
    ".pdf":  "PDF",
    ".docx": "Word Document",
    ".pptx": "PowerPoint Presentation",
    ".xlsx": "Excel Spreadsheet",
    ".xls":  "Excel Spreadsheet",
    ".csv":  "CSV Spreadsheet",
    ".txt":  "Text File",
    ".md":   "Markdown File",
    ".html": "HTML File",
    ".htm":  "HTML File",
    ".jpg":  "Image (JPEG)",
    ".jpeg": "Image (JPEG)",
    ".png":  "Image (PNG)",
    ".bmp":  "Image (BMP)",
    ".tiff": "Image (TIFF)",
    ".gif":  "Image (GIF)",
}


# ---------------------------------------------------------------------------
# Content-signature (magic byte) verification
# ---------------------------------------------------------------------------

# Leading byte signatures keyed by extension. ZIP-based Office formats
# (.docx/.pptx/.xlsx) share the "PK" signature. Plain-text formats
# (.csv/.txt/.md/.html) have no reliable magic and are skipped.
_MAGIC_SIGNATURES: dict[str, tuple[bytes, ...]] = {
    ".pdf":  (b"%PDF",),
    ".docx": (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08"),
    ".pptx": (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08"),
    ".xlsx": (b"PK\x03\x04", b"PK\x05\x06", b"PK\x07\x08"),
    ".xls":  (b"\xd0\xcf\x11\xe0",),               # OLE2 compound document
    ".png":  (b"\x89PNG\r\n\x1a\n",),
    ".jpg":  (b"\xff\xd8\xff",),
    ".jpeg": (b"\xff\xd8\xff",),
    ".gif":  (b"GIF87a", b"GIF89a"),
    ".bmp":  (b"BM",),
    ".tiff": (b"II*\x00", b"MM\x00*"),
}

# Office formats that are really ZIP archives — guard against decompression bombs.
_ZIP_OFFICE_EXTS = {".docx", ".pptx", ".xlsx"}


def _verify_magic(file_path: Path, ext: str) -> None:
    """Verify the file's leading bytes match its claimed extension.

    Defends against a renamed/spoofed file (e.g. an executable named .pdf).
    Extensions without a reliable signature are skipped.
    """
    sigs = _MAGIC_SIGNATURES.get(ext)
    if not sigs:
        return
    with open(file_path, "rb") as fh:
        head = fh.read(8)
    if not any(head.startswith(sig) for sig in sigs):
        raise DocumentExtractionError(
            f"File content does not match its '{ext}' extension. "
            "The file may be corrupted or mislabeled."
        )


def _check_zip_bomb(file_path: Path) -> None:
    """Reject zip-based Office files whose uncompressed size dwarfs the upload
    (a decompression bomb that would exhaust memory/disk on parsing)."""
    import zipfile

    max_bytes = settings.max_file_size_mb * 1024 * 1024
    ratio_cap = settings.max_decompression_ratio
    try:
        with zipfile.ZipFile(file_path) as zf:
            total_uncompressed = sum(i.file_size for i in zf.infolist())
    except zipfile.BadZipFile:
        raise DocumentExtractionError("The file is not a valid Office document (bad archive).")

    compressed = max(file_path.stat().st_size, 1)
    # Absolute ceiling: never allow expansion beyond a few hundred MB,
    # and never beyond the configured ratio relative to the upload size.
    if total_uncompressed > ratio_cap * compressed or total_uncompressed > 50 * max_bytes:
        raise DocumentExtractionError(
            "The document expands to an unreasonable size and was rejected "
            "as a potential decompression bomb."
        )


# ---------------------------------------------------------------------------
# Validation (replaces pdf_utils.validate_pdf)
# ---------------------------------------------------------------------------

def validate_document(file_path: Path, original_filename: str) -> None:
    """
    Validate an uploaded document before processing.
    Checks: extension, size, emptiness, content signature, decompression bombs,
    and basic readability. Raises DocumentExtractionError on failure.
    """
    ext = Path(original_filename).suffix.lower()

    if ext not in SUPPORTED_EXTENSIONS:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise DocumentExtractionError(
            f"Unsupported file type '{ext}'. "
            f"Supported formats: {supported}"
        )

    size_mb = file_path.stat().st_size / (1024 * 1024)
    if size_mb > settings.max_file_size_mb:
        raise DocumentExtractionError(
            f"File exceeds the {settings.max_file_size_mb} MB limit "
            f"({size_mb:.1f} MB uploaded)."
        )

    if file_path.stat().st_size == 0:
        raise DocumentExtractionError("The uploaded file is empty.")

    # Content must match the claimed extension (magic-byte check).
    _verify_magic(file_path, ext)

    # Decompression-bomb guard for zip-based Office formats.
    if ext in _ZIP_OFFICE_EXTS:
        _check_zip_bomb(file_path)

    # Format-specific quick check
    if ext == ".pdf":
        try:
            import fitz
            doc = fitz.open(str(file_path))
            if doc.page_count == 0:
                raise DocumentExtractionError("The PDF contains no pages.")
            doc.close()
        except fitz.FileDataError:
            raise DocumentExtractionError("The file appears to be a corrupted PDF.")
    elif ext == ".docx":
        try:
            import docx as _docx
            _docx.Document(str(file_path))
        except Exception:
            raise DocumentExtractionError("Could not open the Word document. The file may be corrupted.")
    elif ext == ".pptx":
        try:
            from pptx import Presentation
            Presentation(str(file_path))
        except Exception:
            raise DocumentExtractionError("Could not open the PowerPoint file. The file may be corrupted.")
    elif ext in (".xlsx", ".xls"):
        try:
            import openpyxl
            openpyxl.load_workbook(str(file_path), read_only=True)
        except Exception:
            raise DocumentExtractionError("Could not open the Excel file. The file may be corrupted.")


# ---------------------------------------------------------------------------
# Main entry point
# ---------------------------------------------------------------------------

def extract_document(file_path: Path, filename: str) -> DocumentContent:
    """
    Universal extraction router.
    Dispatches to the right extractor based on the file extension.
    Always returns a DocumentContent instance.
    """
    ext = Path(filename).suffix.lower()
    label = FORMAT_LABELS.get(ext, ext)
    logger.info(f"[DocUtils] Extracting {label}: {filename}")

    if ext == ".pdf":
        return extract_pdf(file_path, filename)
    elif ext == ".docx":
        return _extract_docx(file_path, filename)
    elif ext == ".pptx":
        return _extract_pptx(file_path, filename)
    elif ext in (".xlsx", ".xls"):
        return _extract_excel(file_path, filename)
    elif ext == ".csv":
        return _extract_csv(file_path, filename)
    elif ext in (".txt", ".md"):
        return _extract_plaintext(file_path, filename)
    elif ext in (".html", ".htm"):
        return _extract_html(file_path, filename)
    elif ext in (".jpg", ".jpeg", ".png", ".bmp", ".tiff", ".gif"):
        return _extract_image(file_path, filename)
    else:
        raise DocumentExtractionError(f"No extractor available for '{ext}'.")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _file_hash(file_path: Path) -> str:
    return hashlib.sha256(file_path.read_bytes()).hexdigest()


def _paginate_text(text: str, lines_per_page: int = 50) -> list[PageContent]:
    """Split a long text into pages of ~N lines each."""
    lines = text.splitlines()
    pages = []
    for i in range(0, max(len(lines), 1), lines_per_page):
        chunk = "\n".join(lines[i : i + lines_per_page]).strip()
        if chunk:
            pages.append(PageContent(page_number=len(pages) + 1, text=chunk))
    if not pages:
        pages.append(PageContent(page_number=1, text=text.strip()))
    return pages


def _make_result(
    file_path: Path,
    filename: str,
    pages: list[PageContent],
    tables: list[TableData] | None = None,
    method: str = "direct",
    native_metadata: dict | None = None,
) -> DocumentContent:
    full_text = "\n\n".join(p.text for p in pages if p.text.strip())
    return DocumentContent(
        file_path=str(file_path),
        filename=filename,
        page_count=len(pages),
        pages=pages,
        full_text=full_text,
        extraction_method=method,
        is_scanned=False,
        file_hash=_file_hash(file_path),
        native_metadata=native_metadata or {},
        tables=tables or [],
    )


def _table_rows_to_markdown(rows: list[list[str]]) -> str:
    """Convert a list of string rows into a markdown table."""
    if len(rows) < 2:
        return ""
    header = [str(c) for c in rows[0]]
    body = [[str(c) for c in row] for row in rows[1:]]
    lines = ["| " + " | ".join(header) + " |",
             "| " + " | ".join(["---"] * len(header)) + " |"]
    for row in body:
        padded = (row + [""] * len(header))[: len(header)]
        lines.append("| " + " | ".join(padded) + " |")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# DOCX
# ---------------------------------------------------------------------------

def _extract_docx(file_path: Path, filename: str) -> DocumentContent:
    try:
        import docx as _docx
    except ImportError:
        raise DocumentExtractionError(
            "python-docx is required to process Word documents. "
            "Run: pip install python-docx"
        )

    doc = _docx.Document(str(file_path))

    # Collect paragraph text; treat every 30 paragraphs as one "page"
    para_texts = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    pages: list[PageContent] = []
    page_size = 30
    for i in range(0, max(len(para_texts), 1), page_size):
        chunk = "\n".join(para_texts[i : i + page_size])
        if chunk.strip():
            pages.append(PageContent(page_number=len(pages) + 1, text=clean_text(chunk)))

    if not pages:
        pages = [PageContent(page_number=1, text="(No text content found in document)")]

    # Extract tables
    tables: list[TableData] = []
    for t_idx, table in enumerate(doc.tables):
        rows = [[cell.text.strip() for cell in row.cells] for row in table.rows]
        if len(rows) < 2:
            continue
        md = _table_rows_to_markdown(rows)
        if md:
            tables.append(TableData(
                page_number=1,  # DOCX tables don't have page numbers without complex rendering
                caption=f"Table {t_idx + 1}",
                markdown=md,
            ))

    # Core properties metadata
    props = doc.core_properties
    meta = {
        "title": props.title or None,
        "author": props.author or None,
        "subject": props.subject or None,
        "created": str(props.created) if props.created else None,
    }

    logger.info(f"[DocUtils] DOCX: {len(pages)} pages, {len(tables)} tables, {len(para_texts)} paragraphs")
    return _make_result(file_path, filename, pages, tables, method="docx", native_metadata=meta)


# ---------------------------------------------------------------------------
# PPTX
# ---------------------------------------------------------------------------

def _extract_pptx(file_path: Path, filename: str) -> DocumentContent:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError:
        raise DocumentExtractionError(
            "python-pptx is required to process PowerPoint files. "
            "Run: pip install python-pptx"
        )

    prs = Presentation(str(file_path))
    pages: list[PageContent] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        slide_text = clean_text("\n".join(texts))
        # Include slide even if empty (preserves slide count)
        pages.append(PageContent(page_number=slide_num, text=slide_text))

    if not pages:
        pages = [PageContent(page_number=1, text="(No text content found in presentation)")]

    # Core properties
    props = prs.core_properties
    meta = {
        "title": props.title or None,
        "author": props.author or None,
        "subject": props.subject or None,
        "created": str(props.created) if props.created else None,
    }

    total_words = sum(p.word_count for p in pages)
    logger.info(f"[DocUtils] PPTX: {len(pages)} slides, {total_words} words")
    return _make_result(file_path, filename, pages, method="pptx", native_metadata=meta)


# ---------------------------------------------------------------------------
# Excel (.xlsx / .xls)
# ---------------------------------------------------------------------------

def _extract_excel(file_path: Path, filename: str) -> DocumentContent:
    try:
        import openpyxl
    except ImportError:
        raise DocumentExtractionError(
            "openpyxl is required to process Excel files. "
            "Run: pip install openpyxl"
        )

    wb = openpyxl.load_workbook(str(file_path), read_only=True, data_only=True)
    pages: list[PageContent] = []
    tables: list[TableData] = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            str_row = [str(cell) if cell is not None else "" for cell in row]
            if any(c.strip() for c in str_row):  # skip entirely empty rows
                rows.append(str_row)

        if not rows:
            pages.append(PageContent(
                page_number=len(pages) + 1,
                text=f"[Sheet: {sheet_name}] (empty)"
            ))
            continue

        # Page = text representation of the sheet
        text_lines = [f"Sheet: {sheet_name}"]
        for row in rows:
            text_lines.append("\t".join(row))
        page_text = clean_text("\n".join(text_lines))
        pages.append(PageContent(page_number=len(pages) + 1, text=page_text))

        # Also store as a proper markdown table
        if len(rows) >= 2:
            md = _table_rows_to_markdown(rows)
            if md:
                tables.append(TableData(
                    page_number=len(pages),
                    caption=f"Sheet: {sheet_name}",
                    markdown=md,
                ))

    wb.close()

    if not pages:
        pages = [PageContent(page_number=1, text="(No data found in spreadsheet)")]

    logger.info(f"[DocUtils] Excel: {len(wb.sheetnames)} sheets, {len(tables)} tables")
    return _make_result(file_path, filename, pages, tables, method="excel")


# ---------------------------------------------------------------------------
# CSV
# ---------------------------------------------------------------------------

def _extract_csv(file_path: Path, filename: str) -> DocumentContent:
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        raise DocumentExtractionError(f"Could not read CSV file: {e}")

    reader = csv.reader(text.splitlines())
    rows = [row for row in reader if any(cell.strip() for cell in row)]

    if not rows:
        return _make_result(
            file_path, filename,
            [PageContent(page_number=1, text="(Empty CSV file)")],
            method="csv",
        )

    # Text representation (tab-separated for readability)
    text_lines = ["\t".join(row) for row in rows]
    full_text = clean_text("\n".join(text_lines))
    pages = _paginate_text(full_text, lines_per_page=100)

    # Table view
    tables: list[TableData] = []
    if len(rows) >= 2:
        md = _table_rows_to_markdown(rows[:200])  # cap at 200 rows for markdown
        if md:
            tables.append(TableData(
                page_number=1,
                caption=f"{filename} ({len(rows)} rows)",
                markdown=md,
            ))

    logger.info(f"[DocUtils] CSV: {len(rows)} rows, {len(pages)} pages")
    return _make_result(file_path, filename, pages, tables, method="csv")


# ---------------------------------------------------------------------------
# Plain text / Markdown
# ---------------------------------------------------------------------------

def _extract_plaintext(file_path: Path, filename: str) -> DocumentContent:
    try:
        text = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        raise DocumentExtractionError(f"Could not read file: {e}")

    text = clean_text(text)
    if not text.strip():
        pages = [PageContent(page_number=1, text="(Empty file)")]
    else:
        pages = _paginate_text(text, lines_per_page=50)

    ext = Path(filename).suffix.lower()
    method = "markdown" if ext == ".md" else "plaintext"
    logger.info(f"[DocUtils] {method.title()}: {len(pages)} pages")
    return _make_result(file_path, filename, pages, method=method)


# ---------------------------------------------------------------------------
# HTML
# ---------------------------------------------------------------------------

def _extract_html(file_path: Path, filename: str) -> DocumentContent:
    try:
        from bs4 import BeautifulSoup
    except ImportError:
        raise DocumentExtractionError(
            "beautifulsoup4 is required to process HTML files. "
            "Run: pip install beautifulsoup4"
        )

    try:
        raw = file_path.read_text(encoding="utf-8", errors="replace")
    except Exception as e:
        raise DocumentExtractionError(f"Could not read HTML file: {e}")

    soup = BeautifulSoup(raw, "html.parser")

    # Extract page title as metadata
    title_tag = soup.find("title")
    meta_title = title_tag.get_text(strip=True) if title_tag else None

    # Remove script/style noise
    for tag in soup(["script", "style", "noscript", "head"]):
        tag.decompose()

    text = clean_text(soup.get_text(separator="\n"))
    pages = _paginate_text(text, lines_per_page=60) if text.strip() else \
            [PageContent(page_number=1, text="(No readable text found in HTML)")]

    logger.info(f"[DocUtils] HTML: {len(pages)} pages")
    return _make_result(
        file_path, filename, pages, method="html",
        native_metadata={"html_title": meta_title},
    )


# ---------------------------------------------------------------------------
# Image (OCR)
# ---------------------------------------------------------------------------

def _extract_image(file_path: Path, filename: str) -> DocumentContent:
    try:
        import pytesseract
        from PIL import Image
    except ImportError:
        raise DocumentExtractionError(
            "pytesseract and Pillow are required to process image files. "
            "Run: pip install pytesseract Pillow  (also install Tesseract binary)"
        )

    try:
        img = Image.open(str(file_path))
    except Exception as e:
        raise DocumentExtractionError(f"Could not open image: {e}")

    try:
        text = pytesseract.image_to_string(img)
        text = clean_text(text)
    except Exception as e:
        raise DocumentExtractionError(f"OCR failed: {e}")

    if not text.strip():
        pages = [PageContent(page_number=1, text="(No text detected in image)")]
    else:
        pages = _paginate_text(text, lines_per_page=50)

    meta = {
        "image_format": getattr(img, "format", None),
        "image_size": f"{img.width}x{img.height}" if hasattr(img, "width") else None,
        "image_mode": getattr(img, "mode", None),
    }

    logger.info(f"[DocUtils] Image OCR: {sum(p.word_count for p in pages)} words")
    result = _make_result(file_path, filename, pages, method="ocr", native_metadata=meta)
    result.is_scanned = True
    return result
